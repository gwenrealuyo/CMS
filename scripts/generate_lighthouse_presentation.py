#!/usr/bin/env python3
"""Generate The Lighthouse client presentation (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUTPUT = "/Users/gwenvrhernandez/GoCodeo/CMS/docs/The_Lighthouse_Presentation.pptx"

# Harbor Light palette
NAVY = RGBColor(0x1B, 0x2A, 0x41)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
IVORY = RGBColor(0xF7, 0xF4, 0xEF)
CHARCOAL = RGBColor(0x2C, 0x2C, 0x2C)
SLATE = RGBColor(0x5C, 0x66, 0x70)


def _set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _style_title(shape, size=32, color=NAVY):
    shape.text_frame.paragraphs[0].font.size = Pt(size)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.bold = True


def add_logo_placeholder(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(0.2), Inches(1.1), Inches(0.6))
    shape.text = "Logo"
    _set_fill(shape, IVORY)
    shape.line.color.rgb = GOLD
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = NAVY


def add_title_slide(prs, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "The Lighthouse"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "LAMP Church Management System\n"
        "A soul kept is a soul won.\n"
        "Caring for every soul God has placed in our care."
    )
    for i, p in enumerate(subtitle.text_frame.paragraphs):
        p.font.size = Pt(22 if i == 0 else 18)
        p.font.color.rgb = SLATE if i > 0 else CHARCOAL
        if i == 1:
            p.font.bold = True
            p.font.color.rgb = GOLD
    _style_title(slide.shapes.title, size=40)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_logo_placeholder(slide)
    return slide


def add_bullets_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = CHARCOAL
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_logo_placeholder(slide)
    return slide


def add_quote_slide(prs, title, quote, reference, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = f'"{quote}"'
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = reference
    p.font.size = Pt(18)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_logo_placeholder(slide)
    return slide


def add_vision_slide(prs, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Long-Term Vision"
    _style_title(slide.shapes.title)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        ("As God adds to the church, The Lighthouse helps us keep what He wins.", 26, NAVY, True),
        ("", 8, CHARCOAL, False),
        ("Vision: reach and faithfully care for 10,000 souls.", 22, CHARCOAL, False),
        ("Built to grow with the church—so no one is lost in the crowd.", 18, SLATE, False),
    ]
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_logo_placeholder(slide)
    return slide


def add_placeholder_image_slide(prs, title, placeholder_text, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(8.0), Inches(4.5))
    shape.text = placeholder_text
    _set_fill(shape, IVORY)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = SLATE
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_logo_placeholder(slide)
    return slide


def add_architecture_slide(prs, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Architecture Overview"
    _style_title(slide.shapes.title)

    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.8)
    )
    container.text = "On-Premises Church Server"
    _set_fill(container, IVORY)
    container.line.color.rgb = NAVY

    blocks = [
        (Inches(1.4), Inches(2.4), "Frontend\nNext.js"),
        (Inches(4.0), Inches(2.4), "Backend\nDjango + DRF"),
        (Inches(6.6), Inches(2.4), "Database\nSQLite (dev)"),
    ]
    for left, top, label in blocks:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.3), Inches(1.2))
        box.text = label
        _set_fill(box, RGBColor(0xFF, 0xFF, 0xFF))
        box.line.color.rgb = NAVY

    user_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(4.2), Inches(2.3), Inches(0.9)
    )
    user_box.text = "Users\n(Pastor, Deacons, Coordinators)"
    _set_fill(user_box, RGBColor(0xFF, 0xFF, 0xFF))
    user_box.line.color.rgb = NAVY

    slide.shapes.add_connector(1, Inches(3.7), Inches(3.0), Inches(4.0), Inches(3.0))
    slide.shapes.add_connector(1, Inches(6.3), Inches(3.0), Inches(6.6), Inches(3.0))
    slide.shapes.add_connector(1, Inches(2.55), Inches(4.2), Inches(2.55), Inches(3.6))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_logo_placeholder(slide)
    return slide


def main():
    prs = Presentation()

    add_title_slide(
        prs,
        notes="Open with The Lighthouse name, tagline, and pastoral heart. "
        "Tagline: A soul kept is a soul won.",
    )

    add_quote_slide(
        prs,
        "Key Verse",
        "Be thou diligent to know the state of thy flocks, and look well to thy herds.",
        "Proverbs 27:23 (KJV)",
        notes="Anchor stewardship and care, not surveillance.",
    )

    add_vision_slide(
        prs,
        notes="10K souls is a faith vision for scale—same care at every size. "
        "If asked about tracking, use Hebrews 13:17 in conversation only.",
    )

    add_bullets_slide(
        prs,
        "Agenda",
        [
            "Why this matters (church care challenges)",
            "The Lighthouse overview",
            "Usability and daily workflow",
            "Security and on-premises control",
            "Architecture and operations",
            "Roadmap and next steps",
        ],
        notes="15-20 minute comprehensive walkthrough.",
    )

    add_bullets_slide(
        prs,
        "The Problem We Are Solving",
        [
            "Visitors fall through the cracks after first contact",
            "Member data is scattered across spreadsheets or memory",
            "Leaders lack visibility into care, follow-up, and engagement",
            "Ministry coordination becomes reactive instead of proactive",
        ],
        notes="Connect to real ministry gaps.",
    )

    add_bullets_slide(
        prs,
        "The Lighthouse in One Sentence",
        [
            "A ministry care system that helps shepherd every person from first visit to faithful service.",
            "Tagline: A soul kept is a soul won.",
            "Designed for pastors, deacons, and coordinators to work together with clarity.",
        ],
        notes="Keep crisp; repeat tagline once here.",
    )

    add_bullets_slide(
        prs,
        "Usability Focus: The People Journey",
        [
            "Simple intake for visitors and new members",
            "Clear status tracking from visitor to member to servant",
            "Quick search and filters to find people fast",
            "Role-based navigation so each leader sees what they need",
        ],
        notes="Emphasize everyday ease of use.",
    )

    add_bullets_slide(
        prs,
        "Core Modules at a Glance",
        [
            "People and Families management",
            "Clusters and attendance tracking",
            "Events and lessons",
            "Ministries and evangelism",
            "Finance hub (admin/pastor access)",
        ],
        notes="Quick orientation before screenshots.",
    )

    add_placeholder_image_slide(
        prs,
        "Usability Highlights (Screenshots)",
        "Placeholder: People dashboard, search, and filters",
        notes="Insert UI screenshots.",
    )

    add_bullets_slide(
        prs,
        "Security You Control (On-Premises)",
        [
            "Runs on your own server - no third-party cloud required",
            "JWT authentication with short-lived access tokens",
            "Role-based access control for pastors, coordinators, members",
            "Branch-based data filtering for multi-branch churches",
        ],
        notes="Emphasize control, privacy, stewardship.",
    )

    add_bullets_slide(
        prs,
        "Access Control in Practice",
        [
            "ADMIN/PASTOR: full access",
            "COORDINATOR: module-specific access",
            "MEMBER: read-only for their own data",
            "VISITOR: cannot log in",
        ],
        notes="Least-privilege and safety.",
    )

    add_architecture_slide(
        prs,
        notes="Frontend -> backend -> database, all on church server.",
    )

    add_bullets_slide(
        prs,
        "Operations and Setup",
        [
            "Local server setup with Docker-supported database",
            "Admin user creation and role assignment",
            "Daily use through a web browser",
            "Backups managed by church IT or designated admin",
        ],
        notes="Reassure non-technical leaders.",
    )

    add_bullets_slide(
        prs,
        "Roadmap: Future Enhancements",
        [
            "Email-based password reset",
            "Optional token blacklisting",
            "Two-factor authentication for sensitive roles",
            "Session management and audit visibility",
        ],
        notes="Forward-looking without overpromising.",
    )

    add_placeholder_image_slide(
        prs,
        "Demo or Walkthrough",
        "Placeholder: Live demo or recorded walkthrough",
        notes="Live demo or recording.",
    )

    add_bullets_slide(
        prs,
        "Next Steps",
        [
            "Confirm priorities with pastors and coordinators",
            "Finalize branding and ministry workflows",
            "Load initial data and train leaders",
            "Set timeline for launch",
        ],
        notes="Clear call to action.",
    )

    add_bullets_slide(
        prs,
        "Q&A",
        ["Questions and feedback"],
        notes="Invite discussion.",
    )

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")


if __name__ == "__main__":
    main()
